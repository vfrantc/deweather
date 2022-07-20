import os
import math
import cv2
import json
from scipy.io import loadmat
from pptx import Presentation
from pptx.util import Inches


def form_grid(input_dir, methods, imname, num_cols = 4, rows_per_page = 3):
    num_pages = int(math.ceil(len(methods) / (num_cols * rows_per_page)))

    i = 0
    pages = []
    for page_num in range(num_pages):
        row = []
        pages.append(row)
        for row_num in range(rows_per_page):
            column = []
            row.append(column)
            for column_num in range(num_cols):
                if i == len(methods):
                    break

                column.append((methods[i], os.path.join(input_dir, methods[i], imname)))

                i += 1

    return pages


def insert_tables(prs, pres_dataset = 'rain12'):
    input_filename = '../benchmark/results_real.mat'
    title_only_slide_layout = prs.slide_layouts[5]
    mat = loadmat(input_filename)

    # form tables out of matrices
    tables = []
    for i, dataset in enumerate(mat['datasets'][0]):
        if dataset[0] != pres_dataset:
            continue

        myds = dataset[0]
        table = [['']]

        for method in mat['metrics'][0]:
            table[0].append(method[0].upper())

        for j, method in enumerate(mat['methods'][0]):
            table.append([method[0]])
            for k, metric in enumerate(mat['metrics'][0]):
                table[j+1].append(mat['results'][i,j,k])

        tables.append(table)

    # for each table (each dataset)
    for mtable in tables:

        slide = prs.slides.add_slide(title_only_slide_layout)
        slide.shapes.title.text = myds
        rows = len(mtable)
        cols = len(mtable[0])
        top = Inches(1.5)
        left = Inches(0.0)
        width = Inches(10.0)
        height = Inches(6.0)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for i, row in enumerate(mtable):
            for j, val in enumerate(row):
                if isinstance(val, float):
                    table.cell(i, j).text = '{:.3f}'.format(val)
                else:
                    table.cell(i, j).text = str(val)


def insert_images(prs, folder, methods, files, num_cols = 3, num_rows = 3):
    blank_slide_layout = prs.slide_layouts[6]

    for imname in files:
        pages = form_grid(folder, methods, imname, num_cols = num_cols, rows_per_page = num_rows)
        print(os.path.join(folder, 'input', imname))
        img = cv2.imread(os.path.join(folder, 'input', imname))
        im_height, im_width = img.shape[:2]
        ratio = im_width / im_height
        width_inch = 3.2
        height_inch = width_inch / ratio

        for page in pages:
            slide = prs.slides.add_slide(blank_slide_layout)
            for n_row, row in enumerate(page):
                for n_col, col in enumerate(row):
                    label, img_path = col
                    left = Inches(0.1 + (width_inch + 0.1)*n_col)
                    top = Inches(0.5 + (0.32 + height_inch)*n_row)

                    text_left = Inches(0.1 + (width_inch + 0.1) * n_col)
                    text_top = Inches(0.1 + (height_inch + 0.32) * (n_row+1))

                    pic = slide.shapes.add_picture(img_path, left, top, Inches(width_inch), Inches(height_inch))
                    txBox = slide.shapes.add_textbox(text_left, text_top, Inches(width_inch), Inches(0.1))
                    tf = txBox.text_frame
                    tf.text = label.upper()

def title_only_slide(prs, title='Real-world images'):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title

if __name__ == '__main__':
    '''
    prs = Presentation()

    descr = {'Method': 'Paper',
             'GCA': 'Chen, Dongdong, et al. "Gated context aggregation network for image dehazing and deraining." 2019 IEEE winter conference on applications of computer vision (WACV). IEEE, 2019.',
             'GCN': 'Fu, X., Qi, Q., Zha, Z. J., Zhu, Y., & Ding, X. (2021, May). Rain streak removal via dual graph convolutional network. In Proc. AAAI Conf. Artif. Intell. (pp. 1-9).',
             'IRR': 'Wei, W., Meng, D., Zhao, Q., Xu, Z., & Wu, Y. (2019). Semi-supervised transfer learning for image rain removal. In Proceedings of the IEEE/CVF Conference on Computer Vision and Pattern Recognition (pp. 3877-3886).',
             'MPRNET': 'Zamir, S. W., Arora, A., Khan, S., Hayat, M., Khan, F. S., Yang, M. H., & Shao, L. (2021). Multi-stage progressive image restoration. In Proceedings of the IEEE/CVF Conference on Computer Vision and Pattern Recognition (pp. 14821-14831).',
             'VRG': 'Wang, H., Yue, Z., Xie, Q., Zhao, Q., Zheng, Y., & Meng, D. (2021). From Rain Generation to Rain Removal. In Proceedings of the IEEE/CVF Conference on Computer Vision and Pattern Recognition (pp. 14791-14801).'}

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Methods'

    rows , cols = len(descr), 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)

    # write column headings
    for i, (name, des) in enumerate(descr.items()):
        table.cell(i, 0).text = name
        table.cell(i, 1).text = des


    with open('images.json', 'r') as f:
        data = json.load(f)

    for name, dataset in data.items():

        if name == 'real-world-images':
            prs.save('synthetic.pptx')
            prs = Presentation()

        if name == 'rain12600':
            out = []
            for imname in dataset['images']:
                for i in range(1,5):
                    #out.append(imname)
                    out.append(imname.split('.')[0] + '_{}.png'.format(i))
            images = out
        else:
            images = dataset['images']

        if name in ['rain12', 'rain200H', 'rain200L', 'rain12600']:
            insert_tables(prs, pres_dataset=name)
        else:
            title_only_slide(prs, title=name)
        insert_images(prs, dataset['processed_path'].replace('out', 'out2'), dataset['methods'], images)

    prs.save('real-world-images.pptx')
    '''

    '''
    prs = Presentation()

    with open('images.json', 'r') as f:
        data = json.load(f)

    for name, dataset in data.items():
        if name in ['rain12', 'rain200H', 'rain200L', 'rain12600']:
            insert_tables(prs, pres_dataset=name)

    prs.save('real-world-images.pptx')

'''
    prs = Presentation()
    insert_tables(prs, pres_dataset='Test100')

    with open('images.json', 'r') as f:
        data = json.load(f)
    for name, dataset in data.items():
        insert_tables(prs, pres_dataset=name)
        title_only_slide(prs, title=name)

        images = dataset['images']
        insert_images(prs, dataset['processed_path'].replace('out', 'out2'), dataset['methods'], images)

    prs.save('out2.pptx')